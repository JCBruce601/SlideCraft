"""
SlideCraft v5.0 - Streamlit Web Application
Presentations that inspire, in seconds
"""

import streamlit as st
from slidecraft_v5 import (
    create_presentation,
    ThemeGallery,
    TEMPLATE_LIBRARY,
    BrandKit
)
from ai_generator import generate_with_ai
import os
from datetime import datetime
from dotenv import load_dotenv
from pathlib import Path

# Load environment variables from .env file in the same directory as this script
env_path = Path(__file__).parent / '.env'
load_dotenv(dotenv_path=env_path)

# Page config
st.set_page_config(
    page_title="SlideCraft v5.0",
    page_icon="🎨",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize session state for dark mode
if 'dark_mode' not in st.session_state:
    st.session_state.dark_mode = False

# Custom CSS with dark mode support
if st.session_state.dark_mode:
    # Dark mode colors
    bg_color = "#1a1a1a"
    text_color = "#e0e0e0"
    tagline_color = "#b0b0b0"
    card_border = "#404040"
    card_hover = "#667eea"
    card_bg = "#2a2a2a"
else:
    # Light mode colors
    bg_color = "#ffffff"
    text_color = "#000000"
    tagline_color = "#666666"
    card_border = "#e0e0e0"
    card_hover = "#667eea"
    card_bg = "#ffffff"

st.markdown(f"""
<style>
    /* Main background */
    .stApp {{
        background-color: {bg_color};
        color: {text_color};
    }}

    /* Sidebar */
    [data-testid="stSidebar"] {{
        background-color: {bg_color};
    }}

    .main-header {{
        font-size: 3rem;
        font-weight: bold;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.5rem;
    }}
    .tagline {{
        font-size: 1.2rem;
        color: {tagline_color};
        margin-bottom: 2rem;
    }}
    .theme-card {{
        padding: 1rem;
        border-radius: 8px;
        border: 2px solid {card_border};
        margin: 0.5rem 0;
        background-color: {card_bg};
    }}
    .theme-card:hover {{
        border-color: {card_hover};
        box-shadow: 0 2px 8px rgba(102, 126, 234, 0.2);
    }}

    /* Input fields in dark mode */
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea,
    .stNumberInput > div > div > input {{
        background-color: {card_bg};
        color: {text_color};
        border-color: {card_border};
    }}

    /* Selectbox in dark mode */
    .stSelectbox > div > div {{
        background-color: {card_bg};
        color: {text_color};
    }}
</style>
""", unsafe_allow_html=True)

# Header
st.markdown('<div class="main-header">🎨 SlideCraft</div>', unsafe_allow_html=True)
st.markdown('<div class="tagline">Presentations that inspire, in seconds</div>', unsafe_allow_html=True)

# Sidebar
with st.sidebar:
    st.header("⚙️ Settings")

    # Dark mode toggle
    col1, col2 = st.columns([3, 1])
    with col1:
        st.write("**Theme**")
    with col2:
        if st.button("🌙" if not st.session_state.dark_mode else "☀️", key="dark_mode_toggle"):
            st.session_state.dark_mode = not st.session_state.dark_mode
            st.rerun()

    st.divider()

    # Mode selection
    mode = st.radio(
        "Creation Mode",
        ["AI-Powered ✨", "Quick Create", "Template-Based", "Custom Template Upload"],
        help="AI-Powered: Paste notes/agenda and let Claude build everything. Quick Create for custom topics. Template-Based for structured presentations. Custom Template Upload for your own .pptx templates."
    )

    st.divider()
    
    # Theme selection
    st.subheader("🎨 Theme")
    themes = ThemeGallery.list_themes()
    theme_names = {
        'software_professional': '💻 Software Professional',
        'church_warmth': '⛪ Church Warmth',
        'startup_vibrant': '🚀 Startup Vibrant',
        'executive_minimal': '💼 Executive Minimal',
        'creative_bold': '🎭 Creative Bold',
        'tech_modern': '💻 Tech Modern',
        'healthcare_trust': '🏥 Healthcare Trust',
        'education_friendly': '📚 Education Friendly',
        'finance_corporate': '💰 Finance Corporate',
        'marketing_dynamic': '📢 Marketing Dynamic',
        'nonprofit_warm': '🤝 Nonprofit Warm',
    }
    
    selected_theme = st.selectbox(
        "Select Theme",
        themes,
        format_func=lambda x: theme_names.get(x, x)
    )
    
    # Show theme preview
    theme = ThemeGallery.get_theme(selected_theme)
    st.caption(f"**{theme.name}**")
    st.caption(theme.description)
    
    # Show theme colors
    cols = st.columns(5)
    color_keys = ['primary', 'primary_light', 'secondary', 'accent', 'accent_light']
    for i, key in enumerate(color_keys):
        with cols[i]:
            color = theme.colors.get(key, '#000000')
            st.markdown(f'<div style="background-color: {color}; height: 30px; border-radius: 4px; border: 1px solid #ccc;"></div>', unsafe_allow_html=True)

# Main content area
if mode == "AI-Powered ✨":
    st.header("✨ AI-Powered Generation")
    st.write("Just provide your topic and paste your content - Claude will create a professional presentation for you.")

    col1, col2 = st.columns([2, 1])

    with col1:
        topic = st.text_input(
            "Presentation Topic *",
            placeholder="e.g., Walking in Faith, Q4 Business Review, Introduction to Machine Learning..."
        )

        content = st.text_area(
            "Your Content (paste anything) *",
            placeholder="""Paste your notes, meeting agenda, sermon outline, scripture references, bullet points, etc.

Example for sermon:
Scripture: Hebrews 11:1-6

Main Points:
1. Faith is substance - it's real and tangible
2. Faith pleases God - without it we can't please Him
3. Faith requires action - examples from Hall of Faith

Application:
- Trust God in one specific area this week
- Look for evidence of God's faithfulness

Example for business:
- Revenue up 23% YoY
- New product launch exceeded targets
- Customer retention improved to 87%
- Q1 goals: expand into new markets, hire 5 engineers
            """,
            height=300
        )

    with col2:
        presentation_type = st.selectbox(
            "Presentation Type *",
            ["general", "sermon", "business", "education"],
            format_func=lambda x: {
                'general': '📊 General',
                'sermon': '⛪ Sermon',
                'business': '💼 Business',
                'education': '📚 Education'
            }[x],
            help="Type helps Claude structure the presentation appropriately"
        )

        num_slides_ai = st.number_input(
            "Target Number of Slides (optional)",
            min_value=0,
            max_value=30,
            value=0,
            help="Leave at 0 to let Claude decide the optimal number"
        )

        # API Key input (check for .env first)
        api_key_from_env = os.getenv('ANTHROPIC_API_KEY')

        if api_key_from_env:
            st.success("✅ API key loaded from .env file")
            st.caption(f"🔑 Key: {api_key_from_env[:20]}... (hidden)")
            api_key = api_key_from_env
        else:
            st.warning("⚠️ No API key found in .env file")
            st.caption("💡 If you just created .env, restart Streamlit with Ctrl+C then `streamlit run app.py`")
            api_key = st.text_input(
                "Anthropic API Key *",
                type="password",
                placeholder="sk-ant-...",
                help="Get your API key from https://console.anthropic.com/"
            )

    st.divider()

    if st.button("✨ Generate with AI", type="primary", use_container_width=True):
        if not topic:
            st.error("⚠️ Please enter a presentation topic")
        elif not content:
            st.error("⚠️ Please provide your content (notes, agenda, outline, etc.)")
        elif not api_key:
            st.error("⚠️ Please enter your Anthropic API key or add it to .env file")
        else:
            with st.spinner("🤖 Claude is analyzing your content and creating slides..."):
                try:
                    # Generate presentation with AI
                    config = generate_with_ai(
                        topic=topic,
                        content=content,
                        theme=selected_theme,
                        num_slides=num_slides_ai if num_slides_ai > 0 else None,
                        presentation_type=presentation_type,
                        api_key=api_key
                    )

                    # Create the presentation
                    result = create_presentation(config)

                    st.success("✅ AI-powered presentation created successfully!")

                    # Show what was generated
                    st.info(f"📊 Generated {len(config['slides_content'])} slides with speaker notes")

                    # Download button
                    if os.path.exists(result['filepath']):
                        with open(result['filepath'], 'rb') as f:
                            st.download_button(
                                label="📥 Download AI-Generated Presentation",
                                data=f,
                                file_name=f"slidecraft_ai_{topic.lower().replace(' ', '_')[:30]}_{datetime.now().strftime('%Y%m%d')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                type="primary"
                            )

                        # Show preview of generated structure
                        with st.expander("📋 View Generated Structure"):
                            for i, slide in enumerate(config['slides_content'], 1):
                                st.write(f"**Slide {i}:** {slide.get('type', 'content').title()} - {slide.get('title', 'N/A')}")

                        st.info(f"💡 **Pro Tip:** The AI added comprehensive speaker notes to help you present. Open in PowerPoint to view them.")
                    else:
                        st.error("Presentation file not found. Please try again.")

                except ValueError as ve:
                    st.error(f"⚠️ Error: {str(ve)}")
                    st.caption("Check that your API key is valid and you have credits available.")
                except Exception as e:
                    st.error(f"❌ Error generating presentation: {str(e)}")
                    st.exception(e)

elif mode == "Quick Create":
    st.header("✨ Quick Create")
    st.write("Describe your presentation and let SlideCraft build it for you.")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        topic = st.text_input(
            "Presentation Topic *",
            placeholder="e.g., Q4 Business Review, Sunday Sermon on Faith, Product Launch..."
        )
        
        context = st.text_area(
            "Additional Context (optional)",
            placeholder="Paste meeting agenda, outline, or key points...",
            height=150
        )
    
    with col2:
        num_slides = st.number_input(
            "Number of Slides",
            min_value=3,
            max_value=30,
            value=8
        )
        
        company_name = st.text_input(
            "Company/Organization (optional)",
            placeholder="Your organization name"
        )
        
        presenter = st.text_input(
            "Presenter Name (optional)",
            placeholder="Your name"
        )
    
    st.divider()
    
    if st.button("🎨 Generate Presentation", type="primary", use_container_width=True):
        if not topic:
            st.error("⚠️ Please enter a presentation topic")
        elif num_slides < 3 or num_slides > 30:
            st.error("⚠️ Number of slides must be between 3 and 30")
        else:
            with st.spinner("Creating your presentation..."):
                try:
                    # Validate topic length
                    if len(topic) > 200:
                        st.warning("Topic is very long. Truncating to 200 characters.")
                        topic = topic[:200]

                    # Build slides content
                    slides_content = [
                        {'type': 'title', 'title': topic, 'subtitle': f'{company_name}\n{presenter}' if company_name or presenter else ''}
                    ]

                    # Add content slides based on context if provided
                    if context:
                        # Parse context intelligently
                        lines = [l.strip() for l in context.split('\n') if l.strip()]

                        # Add agenda slide if we have enough content
                        if len(lines) >= 3:
                            slides_content.append({
                                'type': 'content',
                                'title': 'Overview',
                                'bullets': lines[:min(5, len(lines))]
                            })

                        # Create content slides from remaining lines
                        remaining_lines = lines[5:] if len(lines) > 5 else lines
                        current_slide_bullets = []
                        current_title = None

                        for line in remaining_lines:
                            # Check if line looks like a header (short, ends with :, or is numbered)
                            is_header = (
                                len(line) < 50 and
                                (line.endswith(':') or
                                 line.startswith(tuple(f'{i}.' for i in range(1, 10))) or
                                 line.startswith(tuple(f'{i})' for i in range(1, 10))))
                            )

                            if is_header:
                                # Save previous slide if we have content
                                if current_title and current_slide_bullets:
                                    slides_content.append({
                                        'type': 'content',
                                        'title': current_title,
                                        'bullets': current_slide_bullets[:6]  # Max 6 bullets
                                    })

                                # Start new slide
                                current_title = line.rstrip(':').strip()
                                current_slide_bullets = []
                            else:
                                # Add as bullet point
                                if current_title is None:
                                    current_title = f'Key Points'

                                current_slide_bullets.append(line)

                                # Create slide if we have enough bullets
                                if len(current_slide_bullets) >= 6:
                                    slides_content.append({
                                        'type': 'content',
                                        'title': current_title,
                                        'bullets': current_slide_bullets[:6]
                                    })
                                    current_title = None
                                    current_slide_bullets = []

                        # Add final slide if we have remaining content
                        if current_title and current_slide_bullets:
                            slides_content.append({
                                'type': 'content',
                                'title': current_title,
                                'bullets': current_slide_bullets[:6]
                            })

                        # Fill remaining slides if needed
                        while len(slides_content) < num_slides - 1:
                            slides_content.append({
                                'type': 'content',
                                'title': f'{topic} - Additional Points',
                                'bullets': [
                                    'Key insight',
                                    'Supporting evidence',
                                    'Practical application'
                                ]
                            })
                    else:
                        # Better generic content generation
                        # Add agenda
                        slides_content.append({
                            'type': 'content',
                            'title': 'Agenda',
                            'bullets': [
                                'Introduction and Overview',
                                'Main Discussion Points',
                                'Key Takeaways',
                                'Questions & Next Steps'
                            ]
                        })

                        # Add section breaks and content slides
                        sections = [
                            ('Introduction', ['Overview of topic', 'Background context', 'Objectives for today']),
                            ('Main Points', ['First key concept', 'Supporting details', 'Real-world examples']),
                            ('Deep Dive', ['Technical details', 'Data and analysis', 'Case studies']),
                            ('Key Takeaways', ['Summary of main points', 'Action items', 'Resources and next steps'])
                        ]

                        slides_per_section = (num_slides - 2) // len(sections)

                        for section_title, default_bullets in sections:
                            if len(slides_content) >= num_slides - 1:
                                break

                            # Add section divider
                            slides_content.append({
                                'type': 'section',
                                'title': section_title,
                                'section_number': str(len([s for s in slides_content if s.get('type') == 'section']) + 1)
                            })

                            # Add content slides for this section
                            for i in range(min(slides_per_section, num_slides - len(slides_content) - 1)):
                                slides_content.append({
                                    'type': 'content',
                                    'title': f'{section_title} - Details',
                                    'bullets': default_bullets
                                })

                    # Add closing slide
                    slides_content.append({
                        'type': 'section',
                        'title': 'Thank You',
                        'section_number': ''
                    })

                    # Create presentation
                    config = {
                        'theme': selected_theme,
                        'slides_content': slides_content
                    }

                    result = create_presentation(config)

                    st.success("✅ Presentation created successfully!")

                    # Download button
                    if os.path.exists(result['filepath']):
                        with open(result['filepath'], 'rb') as f:
                            st.download_button(
                                label="📥 Download Presentation",
                                data=f,
                                file_name=f"slidecraft_{topic.lower().replace(' ', '_')[:30]}_{datetime.now().strftime('%Y%m%d')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                type="primary"
                            )

                        st.info(f"💡 **Tip:** Open in PowerPoint or Google Slides. Contains {len(slides_content)} slides with speaker notes.")
                    else:
                        st.error("Presentation file not found. Please try again.")

                except ValueError as ve:
                    st.error(f"⚠️ Validation error: {str(ve)}")
                except IOError as ioe:
                    st.error(f"⚠️ File error: {str(ioe)}. Check that you have write permissions.")
                except Exception as e:
                    st.error(f"❌ Error creating presentation: {str(e)}")
                    st.exception(e)

elif mode == "Template-Based":
    st.header("📋 Template-Based Creation")
    st.write("Choose a professional template and fill in the details.")
    
    # Organize templates by category
    categories = {}
    for tid, template in TEMPLATE_LIBRARY.items():
        cat = template['category']
        if cat not in categories:
            categories[cat] = []
        categories[cat].append((tid, template))
    
    # Category icons
    category_icons = {
        'church': '⛪',
        'business': '💼',
        'marketing': '📢',
        'education': '📚',
        'government': '🏛️'
    }
    
    # Category selection
    selected_category = st.selectbox(
        "Select Category",
        list(categories.keys()),
        format_func=lambda x: f"{category_icons.get(x, '📁')} {x.title()}"
    )
    
    st.divider()
    
    # Template selection
    category_templates = categories[selected_category]
    
    cols = st.columns(len(category_templates))
    selected_template_id = None
    
    for idx, (tid, template) in enumerate(category_templates):
        with cols[idx]:
            if st.button(
                f"**{template['name']}**\n\n{template['description']}", 
                use_container_width=True,
                key=f"template_{tid}"
            ):
                selected_template_id = tid
    
    # Store selected template in session state
    if selected_template_id:
        st.session_state.selected_template = selected_template_id
    
    if 'selected_template' in st.session_state:
        template_id = st.session_state.selected_template
        template = TEMPLATE_LIBRARY[template_id]
        
        st.divider()
        st.subheader(f"📝 {template['name']}")
        st.caption(template['description'])
        
        # Create form for template fields
        template_data = {}
        
        # Get unique fields from template slides
        fields = set()
        for slide in template['slides']:
            # Extract fields from title
            if 'title' in slide:
                import re
                fields.update(re.findall(r'\{([^}]+)\}', slide['title']))
            if 'subtitle' in slide:
                fields.update(re.findall(r'\{([^}]+)\}', slide['subtitle']))
            if 'bullets' in slide:
                for bullet in slide['bullets']:
                    fields.update(re.findall(r'\{([^}]+)\}', bullet))
            if 'left_items' in slide:
                for item in slide['left_items']:
                    fields.update(re.findall(r'\{([^}]+)\}', item))
            if 'right_items' in slide:
                for item in slide['right_items']:
                    fields.update(re.findall(r'\{([^}]+)\}', item))
        
        # Create input fields
        col1, col2 = st.columns(2)
        fields_list = sorted(list(fields))
        
        for idx, field in enumerate(fields_list):
            with col1 if idx % 2 == 0 else col2:
                template_data[field] = st.text_input(
                    field.replace('_', ' ').title(),
                    key=f"field_{field}",
                    placeholder=f"Enter {field.replace('_', ' ')}"
                )
        
        st.divider()
        
        if st.button("🎨 Generate from Template", type="primary", use_container_width=True):
            # Check if required fields are filled
            missing_fields = [f for f, v in template_data.items() if not v]

            if missing_fields:
                st.warning(f"⚠️ Please fill in all fields. Missing: {', '.join([f.replace('_', ' ').title() for f in missing_fields[:3]])}")
            else:
                with st.spinner("Creating your presentation..."):
                    try:
                        # Validate field lengths
                        for field, value in template_data.items():
                            if len(value) > 500:
                                st.warning(f"Field '{field}' is very long. Truncating to 500 characters.")
                                template_data[field] = value[:500]

                        # Populate template
                        slides_content = []
                        for slide in template['slides']:
                            slide_copy = slide.copy()

                            # Replace placeholders in title
                            if 'title' in slide_copy:
                                for field, value in template_data.items():
                                    slide_copy['title'] = slide_copy['title'].replace(f'{{{field}}}', value)

                            # Replace in subtitle
                            if 'subtitle' in slide_copy:
                                for field, value in template_data.items():
                                    slide_copy['subtitle'] = slide_copy['subtitle'].replace(f'{{{field}}}', value)

                            # Replace in bullets
                            if 'bullets' in slide_copy:
                                slide_copy['bullets'] = [
                                    bullet.replace(f'{{{field}}}', value)
                                    for bullet in slide_copy['bullets']
                                    for field, value in [list(template_data.items())[0]]  # Trick to get all items
                                ]
                                # Actually replace all fields
                                new_bullets = []
                                for bullet in slide['bullets']:
                                    for field, value in template_data.items():
                                        bullet = bullet.replace(f'{{{field}}}', value)
                                    new_bullets.append(bullet)
                                slide_copy['bullets'] = new_bullets

                            # Replace in two-column items
                            if 'left_items' in slide_copy:
                                new_left = []
                                for item in slide['left_items']:
                                    for field, value in template_data.items():
                                        item = item.replace(f'{{{field}}}', value)
                                    new_left.append(item)
                                slide_copy['left_items'] = new_left

                            if 'right_items' in slide_copy:
                                new_right = []
                                for item in slide['right_items']:
                                    for field, value in template_data.items():
                                        item = item.replace(f'{{{field}}}', value)
                                    new_right.append(item)
                                slide_copy['right_items'] = new_right

                            slides_content.append(slide_copy)

                        # Create presentation
                        config = {
                            'theme': template.get('theme', selected_theme),
                            'slides_content': slides_content
                        }

                        result = create_presentation(config)

                        st.success("✅ Presentation created successfully!")

                        # Download button
                        if os.path.exists(result['filepath']):
                            with open(result['filepath'], 'rb') as f:
                                st.download_button(
                                    label="📥 Download Presentation",
                                    data=f,
                                    file_name=f"slidecraft_{template_id}_{datetime.now().strftime('%Y%m%d')}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    type="primary"
                                )

                            st.info(f"💡 **Tip:** Contains {len(slides_content)} slides with speaker notes. Theme: {template.get('theme', selected_theme)}")
                        else:
                            st.error("Presentation file not found. Please try again.")

                    except ValueError as ve:
                        st.error(f"⚠️ Validation error: {str(ve)}")
                    except IOError as ioe:
                        st.error(f"⚠️ File error: {str(ioe)}. Check that you have write permissions.")
                    except Exception as e:
                        st.error(f"❌ Error creating presentation: {str(e)}")
                        st.exception(e)

elif mode == "Custom Template Upload":
    st.header("📤 Custom Template Upload")
    st.write("Upload your own PowerPoint template (.pptx) and use it to generate presentations.")

    # Template upload
    uploaded_file = st.file_uploader(
        "Upload Template (.pptx)",
        type=['pptx'],
        help="Upload a PowerPoint template file that you want to use for your presentations"
    )

    if uploaded_file is not None:
        # Save uploaded template
        template_dir = os.path.join(os.getcwd(), 'templates', 'custom')
        os.makedirs(template_dir, exist_ok=True)

        template_path = os.path.join(template_dir, uploaded_file.name)

        # Save the file
        with open(template_path, 'wb') as f:
            f.write(uploaded_file.getbuffer())

        st.success(f"✅ Template '{uploaded_file.name}' uploaded successfully!")

        st.info("""
        **Template uploaded!** You can now:
        1. Use the AI-Powered mode and select a custom theme (coming soon)
        2. Use Quick Create with your custom template
        3. Analyze the template structure below

        **Note:** In the current version, uploaded templates are saved but not yet integrated into the generation process.
        This feature will be fully implemented in the next update.
        """)

        # Show template info
        try:
            from pptx import Presentation
            prs = Presentation(template_path)

            st.subheader("📊 Template Analysis")

            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Slides", len(prs.slides))
            with col2:
                st.metric("Layouts", len(prs.slide_layouts))
            with col3:
                slide_width = prs.slide_width / 914400  # Convert to inches
                slide_height = prs.slide_height / 914400
                st.metric("Size", f"{slide_width:.1f}\" x {slide_height:.1f}\"")

            # List available layouts
            st.subheader("Available Layouts")
            for i, layout in enumerate(prs.slide_layouts):
                st.write(f"{i + 1}. {layout.name}")

            st.caption(f"**Template saved to:** {template_path}")

        except Exception as e:
            st.warning(f"Could not analyze template: {str(e)}")

    else:
        st.info("👆 Upload a .pptx template file to get started")

        st.subheader("How to use custom templates:")
        st.markdown("""
        1. **Upload your template** - Any .pptx file with your branding, colors, and layouts
        2. **Template is analyzed** - SlideCraft will detect available layouts and themes
        3. **Generate presentations** - Use your custom template with AI or manual modes

        **Benefits:**
        - Maintain your brand identity across all presentations
        - Use your organization's official templates
        - Leverage professional designs you've purchased or created
        - Combine custom templates with AI-powered content generation

        **Coming Soon:**
        - Select custom templates from a dropdown
        - Map custom layouts to SlideCraft slide types
        - Extract and use custom color themes
        - Save multiple custom templates
        """)

# Footer
st.divider()
col1, col2, col3 = st.columns([1, 2, 1])
with col2:
    st.markdown(f"""
    <div style="text-align: center; color: {tagline_color}; padding: 1rem;">
        <small>SlideCraft v5.0 • Presentations that inspire, in seconds</small><br>
        <small>11 Themes • 10 Templates • Unlimited Possibilities</small>
    </div>
    """, unsafe_allow_html=True)
