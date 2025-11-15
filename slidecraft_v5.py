"""
SlideCraft v5.0 - Presentations that inspire, in seconds
Built on v4 Pro architecture + Church templates

NEW in v5: Church templates (sermon, board, staff), larger fonts, enhanced visibility
Includes: Brand Kits, Theme Gallery (11 themes), Template Library, Smart layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
from io import BytesIO
import os
import json
from datetime import datetime
from dataclasses import dataclass, field, asdict
from typing import Dict, List, Optional

# ============================================================================
# BRAND KIT
# ============================================================================

@dataclass
class BrandKit:
    """Auto-apply brand to presentations"""
    name: str
    logo_path: Optional[str] = None
    primary_colors: List[str] = field(default_factory=list)
    secondary_colors: List[str] = field(default_factory=list)
    heading_font: str = "Arial Bold"
    body_font: str = "Arial"
    
    def to_palette(self):
        """Convert to color palette"""
        def hex_to_rgb(h):
            h = h.lstrip('#')
            return RGBColor(*[int(h[i:i+2], 16) for i in (0, 2, 4)])
        
        palette = {}
        if len(self.primary_colors) >= 1:
            palette['primary'] = hex_to_rgb(self.primary_colors[0])
        if len(self.primary_colors) >= 2:
            palette['primary_light'] = hex_to_rgb(self.primary_colors[1])
        if len(self.secondary_colors) >= 1:
            palette['secondary'] = hex_to_rgb(self.secondary_colors[0])
        if len(self.secondary_colors) >= 2:
            palette['accent'] = hex_to_rgb(self.secondary_colors[1])
        
        palette.setdefault('accent_light', RGBColor(150, 180, 220))
        palette.setdefault('light', RGBColor(240, 245, 250))
        palette.setdefault('text', RGBColor(40, 40, 40))
        palette.setdefault('text_light', RGBColor(100, 100, 100))
        
        return palette
    
    @classmethod
    def from_file(cls, path):
        with open(path) as f:
            return cls(**json.load(f))
    
    def save(self, path):
        with open(path, 'w') as f:
            json.dump(asdict(self), f, indent=2)

# ============================================================================
# THEME GALLERY - 10+ Professional Themes
# ============================================================================

@dataclass
class Theme:
    name: str
    description: str
    colors: Dict[str, str]
    fonts: Dict[str, str]
    style: str
    
    def to_palette(self):
        def hex_to_rgb(h):
            h = h.lstrip('#')
            return RGBColor(*[int(h[i:i+2], 16) for i in (0, 2, 4)])
        return {k: hex_to_rgb(v) for k, v in self.colors.items()}

class ThemeGallery:
    """10+ ready-to-use themes"""
    
    THEMES = {
        'software_professional': Theme(
            'Software Professional', 'Tech/enterprise',
            {'primary': '#1C3A56', 'primary_light': '#2E5470', 'secondary': '#2E86AB',
             'accent': '#16A8A0', 'accent_light': '#5BC0B8', 'light': '#F4F6F7',
             'text': '#2C2C2C', 'text_light': '#6C6C6C'},
            {'heading': 'Arial Bold', 'body': 'Arial'}, 'professional'
        ),
        'startup_vibrant': Theme(
            'Startup Vibrant', 'Bold and energetic',
            {'primary': '#FF6B35', 'primary_light': '#FF8F66', 'secondary': '#004E89',
             'accent': '#F7B801', 'accent_light': '#FFD60A', 'light': '#F8F9FA',
             'text': '#212529', 'text_light': '#6C757D'},
            {'heading': 'Arial Bold', 'body': 'Arial'}, 'bold'
        ),
        'executive_minimal': Theme(
            'Executive Minimal', 'Clean C-suite',
            {'primary': '#2C3E50', 'primary_light': '#34495E', 'secondary': '#7F8C8D',
             'accent': '#3498DB', 'accent_light': '#5DADE2', 'light': '#ECF0F1',
             'text': '#2C3E50', 'text_light': '#95A5A6'},
            {'heading': 'Calibri Bold', 'body': 'Calibri'}, 'minimal'
        ),
        'creative_bold': Theme(
            'Creative Bold', 'Vibrant design/marketing',
            {'primary': '#5D3A9B', 'primary_light': '#7550B5', 'secondary': '#8E44AD',
             'accent': '#E056FD', 'accent_light': '#EB8FFC', 'light': '#F5F0FA',
             'text': '#2C2C2C', 'text_light': '#7C7C7C'},
            {'heading': 'Arial Black', 'body': 'Arial'}, 'creative'
        ),
        'tech_modern': Theme(
            'Tech Modern', 'Modern SaaS',
            {'primary': '#0A192F', 'primary_light': '#1A2938', 'secondary': '#64FFDA',
             'accent': '#00D9FF', 'accent_light': '#66E6FF', 'light': '#F8FAFC',
             'text': '#0F172A', 'text_light': '#64748B'},
            {'heading': 'Arial Bold', 'body': 'Arial'}, 'modern'
        ),
        'healthcare_trust': Theme(
            'Healthcare Trust', 'Medical/healthcare',
            {'primary': '#0E4C92', 'primary_light': '#1A6BBF', 'secondary': '#1EAEDB',
             'accent': '#3FBFB0', 'accent_light': '#6DD4C7', 'light': '#F0F7FA',
             'text': '#333333', 'text_light': '#777777'},
            {'heading': 'Calibri Bold', 'body': 'Calibri'}, 'professional'
        ),
        'education_friendly': Theme(
            'Education Friendly', 'Warm education',
            {'primary': '#2A7C6F', 'primary_light': '#3D9B8C', 'secondary': '#F4A261',
             'accent': '#E76F51', 'accent_light': '#EF8A6F', 'light': '#F7F3E9',
             'text': '#264653', 'text_light': '#6B7780'},
            {'heading': 'Arial Bold', 'body': 'Arial'}, 'friendly'
        ),
        'finance_corporate': Theme(
            'Finance Corporate', 'Conservative finance',
            {'primary': '#003049', 'primary_light': '#1A4A66', 'secondary': '#669BBC',
             'accent': '#C1121F', 'accent_light': '#D64350', 'light': '#F7F9FA',
             'text': '#212529', 'text_light': '#6C757D'},
            {'heading': 'Times New Roman Bold', 'body': 'Times New Roman'}, 'conservative'
        ),
        'marketing_dynamic': Theme(
            'Marketing Dynamic', 'Eye-catching marketing',
            {'primary': '#D62828', 'primary_light': '#E04848', 'secondary': '#F77F00',
             'accent': '#FCBF49', 'accent_light': '#FFD670', 'light': '#FFF8E7',
             'text': '#2C2416', 'text_light': '#7C6F5B'},
            {'heading': 'Arial Bold', 'body': 'Arial'}, 'dynamic'
        ),
        'nonprofit_warm': Theme(
            'Nonprofit Warm', 'Compassionate nonprofit',
            {'primary': '#4A5859', 'primary_light': '#627375', 'secondary': '#7D9D9C',
             'accent': '#C9A96E', 'accent_light': '#E0C589', 'light': '#F5F2ED',
             'text': '#3A3A3A', 'text_light': '#7A7A7A'},
            {'heading': 'Georgia Bold', 'body': 'Georgia'}, 'warm'
        ),
        'church_warmth': Theme(
            'Church Warmth', 'Ministry and worship',
            {'primary': '#78512D', 'primary_light': '#966D4B', 'secondary': '#B3895A',
             'accent': '#C1996B', 'accent_light': '#E6CAAA', 'light': '#FAF8F5',
             'text': '#322816', 'text_light': '#78643C'},
            {'heading': 'Georgia Bold', 'body': 'Georgia'}, 'warm'
        )
    }
    
    @staticmethod
    def get_theme(name):
        return ThemeGallery.THEMES.get(name, ThemeGallery.THEMES['software_professional'])
    
    @staticmethod
    def list_themes():
        return list(ThemeGallery.THEMES.keys())

# ============================================================================
# LAYOUT CONFIG
# ============================================================================

class LayoutConfig:
    WIDESCREEN_16_9 = {
        'name': '16:9 Widescreen', 'width': Inches(13.333), 'height': Inches(7.5),
        'margins': {'top': Inches(0.8), 'bottom': Inches(0.5), 'left': Inches(0.8), 'right': Inches(0.8)},
        'header_height': Inches(1.0), 'content_padding': Inches(0.4), 'line_spacing': 1.3
    }
    
    @staticmethod
    def get_layout(format_type='16:9'):
        return LayoutConfig.WIDESCREEN_16_9

# ============================================================================
# TEMPLATE MANAGER
# ============================================================================

class TemplateManager:
    def __init__(self, template_path=None):
        self.template_path = template_path
        self.prs = None
        self.available_layouts = {}
    
    def load(self):
        if self.template_path and os.path.exists(self.template_path):
            try:
                self.prs = Presentation(self.template_path)
                self._catalog_layouts()
                print(f"✓ Template: {self.template_path}")
            except Exception as e:
                raise ValueError(f"Template error: {e}")
        else:
            self.prs = Presentation()
        return self.prs
    
    def _catalog_layouts(self):
        for idx, layout in enumerate(self.prs.slide_layouts):
            name = layout.name.lower()
            self.available_layouts[name] = idx
            if 'title' in name and 'content' not in name:
                self.available_layouts['title'] = idx
            elif 'content' in name:
                self.available_layouts['content'] = idx
            elif 'section' in name:
                self.available_layouts['section'] = idx
            elif 'blank' in name:
                self.available_layouts['blank'] = idx
    
    def get_layout(self, slide_type):
        if slide_type in self.available_layouts:
            return self.prs.slide_layouts[self.available_layouts[slide_type]]
        return self.prs.slide_layouts[self.available_layouts.get('blank', 6)]
    
    @staticmethod
    def find_custom_templates():
        """Find custom .pptx templates in uploads directory"""
        custom_templates = []
        uploads_dir = '/mnt/user-data/uploads'
        
        if os.path.exists(uploads_dir):
            for file in os.listdir(uploads_dir):
                if file.endswith(('.pptx', '.ppt')):
                    custom_templates.append({
                        'name': file.replace('.pptx', '').replace('.ppt', ''),
                        'path': os.path.join(uploads_dir, file),
                        'filename': file
                    })
        
        return custom_templates

# ============================================================================
# BASE SLIDE (v4 Architecture)
# ============================================================================

class BaseSlide:
    def __init__(self, prs, palette, layout, fonts=None, brand_kit=None):
        self.prs = prs
        self.palette = palette
        self.layout = layout
        self.fonts = fonts or {'heading': 'Arial Bold', 'body': 'Arial'}
        self.brand_kit = brand_kit
        self.slide = None
    
    def create(self, data, template_manager=None):
        layout = template_manager.get_layout(self.slide_type) if template_manager else self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(layout)
        
        self._add_background()
        self._add_brand_logo()
        self._add_header(data)
        self._add_content(data)
        self._add_decorations()
        self._add_speaker_notes(data)
        
        return self.slide
    
    def _add_background(self):
        bg = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.layout['width'], self.layout['height'])
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()
    
    def _add_brand_logo(self):
        if self.brand_kit and self.brand_kit.logo_path and os.path.exists(self.brand_kit.logo_path):
            try:
                self.slide.shapes.add_picture(
                    self.brand_kit.logo_path,
                    self.layout['width'] - Inches(1.5), Inches(0.2),
                    height=Inches(0.5)
                )
            except:
                pass
    
    def _add_header(self, data):
        if hasattr(self, 'skip_header') and self.skip_header:
            return
        header = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.layout['width'], self.layout['header_height'])
        header.fill.solid()
        header.fill.fore_color.rgb = self.palette['primary']
        header.line.fill.background()
    
    def _add_decorations(self):
        pass
    
    def _add_speaker_notes(self, data):
        if data.get('notes'):
            self.slide.notes_slide.notes_text_frame.text = data['notes']
    
    def _add_content(self, data):
        raise NotImplementedError()
    
    @property
    def slide_type(self):
        return 'blank'

# ============================================================================
# SLIDE TYPES
# ============================================================================

class TitleSlide(BaseSlide):
    slide_type = 'title'
    skip_header = True
    
    def _add_background(self):
        bg_bottom = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.layout['width'], self.layout['height'])
        bg_bottom.fill.solid()
        bg_bottom.fill.fore_color.rgb = self.palette['primary']
        bg_bottom.line.fill.background()
        
        bg_top = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.layout['width'], Inches(3))
        bg_top.fill.solid()
        bg_top.fill.fore_color.rgb = self.palette['primary_light']
        bg_top.fill.transparency = 0.3
        bg_top.line.fill.background()
    
    def _add_content(self, data):
        title_box = self.slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.0))
        tf = title_box.text_frame
        tf.text = data.get('title', 'Presentation Title')
        tf.paragraphs[0].font.name = self.fonts['heading']
        tf.paragraphs[0].font.size = Pt(66)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if data.get('subtitle'):
            sub_box = self.slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.2))
            sf = sub_box.text_frame
            sf.text = data['subtitle']
            sf.paragraphs[0].font.name = self.fonts['body']
            sf.paragraphs[0].font.size = Pt(32)
            sf.paragraphs[0].font.color.rgb = self.palette['accent_light']
            sf.paragraphs[0].alignment = PP_ALIGN.CENTER

class ContentSlide(BaseSlide):
    slide_type = 'content'

    def _add_header(self, data):
        super()._add_header(data)
        title_box = self.slide.shapes.add_textbox(Inches(0.9), Inches(0.15), Inches(11), Inches(0.7))
        tf = title_box.text_frame
        tf.text = data.get('title', 'Slide Title')
        tf.paragraphs[0].font.name = self.fonts['heading']
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    def _add_content(self, data):
        # Content background with better positioning
        bg = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.layout['margins']['left'], self.layout['header_height'] + Inches(0.5),
            self.layout['width'] - self.layout['margins']['left'] - self.layout['margins']['right'], Inches(5.4)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.palette['light']
        bg.line.fill.background()

        # Content text box with better spacing
        content_box = self.slide.shapes.add_textbox(
            self.layout['margins']['left'] + Inches(0.6), self.layout['header_height'] + Inches(0.9),
            self.layout['width'] - self.layout['margins']['left'] - self.layout['margins']['right'] - Inches(1.2), Inches(4.8)
        )
        cf = content_box.text_frame
        cf.word_wrap = True
        cf.vertical_anchor = MSO_ANCHOR.TOP

        bullets = data.get('bullets', [])

        # Adjust font size based on number of bullets for better fit
        if len(bullets) <= 3:
            bullet_size = Pt(32)
            space_after = Pt(24)
        elif len(bullets) <= 5:
            bullet_size = Pt(28)
            space_after = Pt(20)
        else:
            bullet_size = Pt(24)
            space_after = Pt(16)

        for idx, bullet in enumerate(bullets):
            p = cf.paragraphs[0] if idx == 0 else cf.add_paragraph()
            p.text = bullet
            p.font.name = self.fonts['body']
            p.font.size = bullet_size
            p.font.color.rgb = self.palette['text']
            p.space_after = space_after
            p.line_spacing = 1.3
            p.level = 0  # Ensure bullet level is set

            # Add bullet point marker
            if idx > 0 or len(bullets) > 1:
                p.text = f"• {bullet}"

class SectionSlide(BaseSlide):
    slide_type = 'section'
    skip_header = True

    def _add_background(self):
        left = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8), self.layout['height'])
        left.fill.solid()
        left.fill.fore_color.rgb = self.palette['secondary']
        left.line.fill.background()

        right = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), 0, self.layout['width'] - Inches(8), self.layout['height'])
        right.fill.solid()
        right.fill.fore_color.rgb = self.palette['primary']
        right.line.fill.background()

    def _add_content(self, data):
        if data.get('section_number'):
            nb = self.slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(2), Inches(1))
            nb.text_frame.text = str(data['section_number'])
            nb.text_frame.paragraphs[0].font.size = Pt(120)
            nb.text_frame.paragraphs[0].font.bold = True
            nb.text_frame.paragraphs[0].font.color.rgb = self.palette['accent_light']

        tb = self.slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(6), Inches(2))
        tb.text_frame.text = data.get('title', 'Section')
        tb.text_frame.word_wrap = True
        tb.text_frame.paragraphs[0].font.name = self.fonts['heading']
        tb.text_frame.paragraphs[0].font.size = Pt(52)
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

class TwoColumnSlide(BaseSlide):
    slide_type = 'two_column'

    def _add_header(self, data):
        super()._add_header(data)
        title_box = self.slide.shapes.add_textbox(Inches(0.9), Inches(0.15), Inches(11), Inches(0.7))
        tf = title_box.text_frame
        tf.text = data.get('title', 'Slide Title')
        tf.paragraphs[0].font.name = self.fonts['heading']
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    def _add_content(self, data):
        # Left column background
        left_bg = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.layout['margins']['left'], self.layout['header_height'] + Inches(0.4),
            Inches(5.5), Inches(5.6)
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = self.palette['light']
        left_bg.line.fill.background()

        # Right column background
        right_bg = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.layout['margins']['left'] + Inches(5.9), self.layout['header_height'] + Inches(0.4),
            Inches(5.5), Inches(5.6)
        )
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = self.palette['light']
        right_bg.line.fill.background()

        # Left column header
        if data.get('left_header'):
            left_header_box = self.slide.shapes.add_textbox(
                self.layout['margins']['left'] + Inches(0.4), self.layout['header_height'] + Inches(0.6),
                Inches(4.7), Inches(0.4)
            )
            lh = left_header_box.text_frame
            lh.text = data['left_header']
            lh.paragraphs[0].font.name = self.fonts['heading']
            lh.paragraphs[0].font.size = Pt(28)
            lh.paragraphs[0].font.bold = True
            lh.paragraphs[0].font.color.rgb = self.palette['primary']

        # Right column header
        if data.get('right_header'):
            right_header_box = self.slide.shapes.add_textbox(
                self.layout['margins']['left'] + Inches(6.3), self.layout['header_height'] + Inches(0.6),
                Inches(4.7), Inches(0.4)
            )
            rh = right_header_box.text_frame
            rh.text = data['right_header']
            rh.paragraphs[0].font.name = self.fonts['heading']
            rh.paragraphs[0].font.size = Pt(28)
            rh.paragraphs[0].font.bold = True
            rh.paragraphs[0].font.color.rgb = self.palette['primary']

        # Left column content
        left_content_box = self.slide.shapes.add_textbox(
            self.layout['margins']['left'] + Inches(0.4), self.layout['header_height'] + Inches(1.2),
            Inches(4.7), Inches(4.5)
        )
        lcf = left_content_box.text_frame
        lcf.word_wrap = True
        lcf.vertical_anchor = MSO_ANCHOR.TOP

        left_items = data.get('left_items', [])
        for idx, item in enumerate(left_items):
            p = lcf.paragraphs[0] if idx == 0 else lcf.add_paragraph()
            p.text = f"• {item}" if len(left_items) > 1 else item
            p.font.name = self.fonts['body']
            p.font.size = Pt(24)
            p.font.color.rgb = self.palette['text']
            p.space_after = Pt(16)
            p.line_spacing = 1.3

        # Right column content
        right_content_box = self.slide.shapes.add_textbox(
            self.layout['margins']['left'] + Inches(6.3), self.layout['header_height'] + Inches(1.2),
            Inches(4.7), Inches(4.5)
        )
        rcf = right_content_box.text_frame
        rcf.word_wrap = True
        rcf.vertical_anchor = MSO_ANCHOR.TOP

        right_items = data.get('right_items', [])
        for idx, item in enumerate(right_items):
            p = rcf.paragraphs[0] if idx == 0 else rcf.add_paragraph()
            p.text = f"• {item}" if len(right_items) > 1 else item
            p.font.name = self.fonts['body']
            p.font.size = Pt(24)
            p.font.color.rgb = self.palette['text']
            p.space_after = Pt(16)
            p.line_spacing = 1.3

class QuoteSlide(BaseSlide):
    slide_type = 'quote'
    skip_header = True

    def _add_background(self):
        # Gradient background effect with shapes
        bg = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.layout['width'], self.layout['height'])
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.palette['primary']
        bg.line.fill.background()

        # Accent shape
        accent = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.5),
            Inches(2), Inches(2)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.palette['accent']
        accent.fill.transparency = 0.7
        accent.line.fill.background()

    def _add_content(self, data):
        # Large quote mark or icon area
        quote_mark = self.slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.0))
        qm = quote_mark.text_frame
        qm.text = '"'
        qm.paragraphs[0].font.size = Pt(180)
        qm.paragraphs[0].font.bold = True
        qm.paragraphs[0].font.color.rgb = self.palette['accent_light']
        qm.paragraphs[0].font.name = self.fonts['heading']

        # Main quote text
        quote_box = self.slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.3), Inches(3.0))
        qt = quote_box.text_frame
        qt.word_wrap = True
        qt.text = data.get('quote', data.get('title', 'Quote text'))
        qt.paragraphs[0].font.name = self.fonts['heading']
        qt.paragraphs[0].font.size = Pt(48)
        qt.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        qt.paragraphs[0].alignment = PP_ALIGN.LEFT
        qt.paragraphs[0].line_spacing = 1.4

        # Attribution
        if data.get('attribution'):
            attr_box = self.slide.shapes.add_textbox(Inches(2.0), Inches(5.8), Inches(9.3), Inches(0.8))
            attr = attr_box.text_frame
            attr.text = f"— {data['attribution']}"
            attr.paragraphs[0].font.name = self.fonts['body']
            attr.paragraphs[0].font.size = Pt(28)
            attr.paragraphs[0].font.color.rgb = self.palette['accent_light']
            attr.paragraphs[0].font.italic = True

class StatsSlide(BaseSlide):
    slide_type = 'stats'
    skip_header = True

    def _add_background(self):
        bg = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.layout['width'], self.layout['height'])
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()

        # Colored accent bar
        accent_bar = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.layout['width'], Inches(1.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.palette['primary']
        accent_bar.line.fill.background()

    def _add_content(self, data):
        # Title at top
        title_box = self.slide.shapes.add_textbox(Inches(1.0), Inches(0.3), Inches(11.3), Inches(1.0))
        tf = title_box.text_frame
        tf.text = data.get('title', 'Key Statistics')
        tf.paragraphs[0].font.name = self.fonts['heading']
        tf.paragraphs[0].font.size = Pt(52)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Stats items (up to 3)
        stats = data.get('stats', [])
        if not stats and data.get('bullets'):
            # Convert bullets to stats if needed
            stats = data['bullets'][:3]

        num_stats = min(len(stats), 3)
        stat_width = Inches(11.3 / num_stats)

        for idx, stat in enumerate(stats[:3]):
            x_pos = Inches(1.0) + (stat_width * idx)

            # Background card
            card = self.slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos + Inches(0.2), Inches(2.5),
                stat_width - Inches(0.4), Inches(3.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.palette['light']
            card.line.color.rgb = self.palette['accent']
            card.line.width = Pt(3)

            # Parse stat (try to extract number and label)
            parts = stat.split(':', 1) if ':' in stat else [stat, '']
            if len(parts) == 2:
                number_text = parts[0].strip()
                label_text = parts[1].strip()
            else:
                # Try to find first number
                import re
                numbers = re.findall(r'[\d,\.]+[%]?', stat)
                if numbers:
                    number_text = numbers[0]
                    label_text = stat.replace(number_text, '').strip()
                else:
                    number_text = stat[:20]
                    label_text = stat[20:]

            # Big number
            num_box = self.slide.shapes.add_textbox(
                x_pos + Inches(0.3), Inches(3.0),
                stat_width - Inches(0.6), Inches(1.5)
            )
            num_tf = num_box.text_frame
            num_tf.text = number_text
            num_tf.paragraphs[0].font.name = self.fonts['heading']
            num_tf.paragraphs[0].font.size = Pt(72)
            num_tf.paragraphs[0].font.bold = True
            num_tf.paragraphs[0].font.color.rgb = self.palette['primary']
            num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Label
            label_box = self.slide.shapes.add_textbox(
                x_pos + Inches(0.3), Inches(4.7),
                stat_width - Inches(0.6), Inches(1.0)
            )
            label_tf = label_box.text_frame
            label_tf.word_wrap = True
            label_tf.text = label_text
            label_tf.paragraphs[0].font.name = self.fonts['body']
            label_tf.paragraphs[0].font.size = Pt(20)
            label_tf.paragraphs[0].font.color.rgb = self.palette['text']
            label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE FACTORY
# ============================================================================

class SlideFactory:
    SLIDE_TYPES = {
        'title': TitleSlide,
        'content': ContentSlide,
        'section': SectionSlide,
        'two_column': TwoColumnSlide,
        'quote': QuoteSlide,
        'stats': StatsSlide
    }

    @staticmethod
    def create_slide(slide_type, prs, palette, layout, fonts=None, brand_kit=None):
        cls = SlideFactory.SLIDE_TYPES.get(slide_type, ContentSlide)
        return cls(prs, palette, layout, fonts, brand_kit)

# ============================================================================
# PRESENTATION BUILDER PRO
# ============================================================================

class PresentationBuilderPro:
    def __init__(self, config):
        self.config = config
        self.brand_kit = None
        self.theme = None
        self.palette = None
        self.fonts = None
        self.layout = LayoutConfig.get_layout(config.get('format', '16:9'))
        self.template_manager = None
        self.prs = None
        self._initialize_styling()
    
    def _initialize_styling(self):
        # Brand Kit > Theme > Default
        if self.config.get('brand_kit'):
            bk = self.config['brand_kit']
            self.brand_kit = BrandKit.from_file(bk) if isinstance(bk, str) else bk
            self.palette = self.brand_kit.to_palette()
            self.fonts = {'heading': self.brand_kit.heading_font, 'body': self.brand_kit.body_font}
            print(f"✓ Brand Kit: {self.brand_kit.name}")
        
        elif self.config.get('theme'):
            self.theme = ThemeGallery.get_theme(self.config['theme'])
            self.palette = self.theme.to_palette()
            self.fonts = self.theme.fonts
            print(f"✓ Theme: {self.theme.name}")
        
        else:
            # Default Tyler
            self.palette = {
                'primary': RGBColor(0, 59, 92), 'primary_light': RGBColor(0, 80, 120),
                'secondary': RGBColor(0, 102, 161), 'accent': RGBColor(74, 144, 226),
                'accent_light': RGBColor(135, 184, 238), 'light': RGBColor(230, 240, 250),
                'text': RGBColor(40, 40, 40), 'text_light': RGBColor(100, 100, 100)
            }
            self.fonts = {'heading': 'Arial Bold', 'body': 'Arial'}
    
    def build(self):
        print("🎨 Building v4 Pro presentation...")
        
        self._initialize_presentation()
        
        self.prs.slide_width = self.layout['width']
        self.prs.slide_height = self.layout['height']
        print(f"✓ Format: {self.layout['name']}")
        
        slides_data = self.config.get('slides_content', [])
        print(f"📊 Creating {len(slides_data)} slides...")
        
        for idx, slide_data in enumerate(slides_data, 1):
            slide_type = slide_data.get('type', 'content')
            title = slide_data.get('title', 'Untitled')[:50]
            print(f"  {idx}. {slide_type}: {title}")
            self._create_slide(slide_type, slide_data)
        
        output_path = self._save()
        print(f"\n✅ Complete: {output_path}")
        return {'filepath': output_path, 'theme': getattr(self.theme, 'name', 'Custom')}
    
    def _initialize_presentation(self):
        template = self.config.get('template_path')
        if template:
            self.template_manager = TemplateManager(template)
            self.prs = self.template_manager.load()
        else:
            self.prs = Presentation()
    
    def _create_slide(self, slide_type, slide_data):
        slide = SlideFactory.create_slide(slide_type, self.prs, self.palette, self.layout, self.fonts, self.brand_kit)
        slide.create(slide_data, self.template_manager)
    
    def _save(self):
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')

        # Try configured output directory, fallback to local outputs/
        output_dir = self.config.get('output_dir', None)

        if not output_dir:
            # Try cloud path, fallback to local
            if os.path.exists('/mnt/user-data/outputs'):
                output_dir = '/mnt/user-data/outputs'
            else:
                output_dir = os.path.join(os.getcwd(), 'outputs')
                os.makedirs(output_dir, exist_ok=True)

        path = os.path.join(output_dir, f"presentation_v4pro_{ts}.pptx")
        self.prs.save(path)
        return path

# ============================================================================
# API
# ============================================================================

# ============================================================================
# TEMPLATE LIBRARY - Church & Business Templates
# ============================================================================

TEMPLATE_LIBRARY = {
    # CHURCH TEMPLATES
    'sermon': {
        'name': 'Sunday Sermon',
        'category': 'church',
        'description': 'Full sermon with scripture, points, and application',
        'theme': 'church_warmth',
        'slides': [
            {'type': 'title', 'title': '{sermon_title}', 'subtitle': '{scripture_reference}\n{date}'},
            {'type': 'section', 'title': 'Opening Scripture', 'section_number': '01'},
            {'type': 'content', 'title': 'Context & Background', 'bullets': [
                'Historical setting: {context_1}',
                'Cultural significance: {context_2}',
                'Why this matters today: {context_3}'
            ]},
            {'type': 'section', 'title': 'Main Point 1: {point_1_title}', 'section_number': '02'},
            {'type': 'content', 'title': '{point_1_title}', 'bullets': [
                '{point_1_detail_1}',
                '{point_1_detail_2}',
                '{point_1_detail_3}'
            ]},
            {'type': 'section', 'title': 'Main Point 2: {point_2_title}', 'section_number': '03'},
            {'type': 'content', 'title': '{point_2_title}', 'bullets': [
                '{point_2_detail_1}',
                '{point_2_detail_2}',
                '{point_2_detail_3}'
            ]},
            {'type': 'content', 'title': 'Practical Application', 'bullets': [
                'This week: {application_1}',
                'In your relationships: {application_2}',
                'Long-term growth: {application_3}'
            ]},
            {'type': 'section', 'title': 'Closing Prayer', 'section_number': '04'}
        ]
    },
    
    'church_board': {
        'name': 'Church Board Meeting',
        'category': 'church',
        'description': 'Leadership meeting with ministry updates and decisions',
        'theme': 'church_warmth',
        'slides': [
            {'type': 'title', 'title': 'Board Meeting', 'subtitle': '{church_name}\n{meeting_date}'},
            {'type': 'content', 'title': 'Agenda', 'bullets': [
                'Opening prayer & attendance',
                'Ministry updates',
                'Financial report',
                'Old business',
                'New business'
            ]},
            {'type': 'section', 'title': 'Ministry Updates', 'section_number': '01'},
            {'type': 'two_column', 'title': 'Attendance & Engagement',
             'left_header': 'This Month', 'left_items': [
                'Attendance: {attendance}',
                'Visitors: {visitors}',
                'Groups: {groups}'
             ],
             'right_header': 'Previous', 'right_items': [
                'Attendance: {prev_attendance}',
                'Visitors: {prev_visitors}',
                'Groups: {prev_groups}'
             ]},
            {'type': 'content', 'title': 'Financial Summary', 'bullets': [
                'Income YTD: {income}',
                'Expenses YTD: {expenses}',
                'Budget status: {budget_status}'
            ]},
            {'type': 'content', 'title': 'Decisions & Action Items', 'bullets': [
                '{decision_1}',
                '{action_item_1}',
                '{action_item_2}'
            ]}
        ]
    },
    
    'staff_meeting': {
        'name': 'Church Staff Meeting',
        'category': 'church',
        'description': 'Weekly staff coordination and prayer',
        'theme': 'church_warmth',
        'slides': [
            {'type': 'title', 'title': 'Staff Meeting', 'subtitle': '{church_name}\n{week_of}'},
            {'type': 'content', 'title': 'This Week', 'bullets': [
                'Sunday services: {services}',
                'Programs: {programs}',
                'Events: {events}'
            ]},
            {'type': 'two_column', 'title': 'Sunday Services',
             'left_header': 'This Sunday', 'left_items': [
                'Theme: {theme}',
                'Worship: {worship}',
                'Tech: {tech}'
             ],
             'right_header': 'Next Sunday', 'right_items': [
                'Theme: {next_theme}',
                'Worship: {next_worship}',
                'Tech: {next_tech}'
             ]},
            {'type': 'content', 'title': 'Prayer Requests', 'bullets': [
                '{prayer_1}',
                '{prayer_2}',
                '{prayer_3}'
            ]}
        ]
    },
    
    # BUSINESS TEMPLATES
    'quarterly_review': {
        'name': 'Quarterly Business Review',
        'category': 'business',
        'description': 'Executive QBR with metrics and highlights',
        'theme': 'software_professional',
        'slides': [
            {'type': 'title', 'title': '{quarter} Business Review', 'subtitle': '{company_name}\n{year}'},
            {'type': 'content', 'title': 'Agenda', 'bullets': [
                'Executive Summary', 'Key Metrics', 'Highlights', 'Priorities'
            ]},
            {'type': 'content', 'title': 'Key Metrics', 'bullets': [
                'Revenue: {revenue}',
                'Growth: {growth}%',
                'Customers: {customers}'
            ]},
            {'type': 'two_column', 'title': 'Performance',
             'left_header': 'Wins', 'left_items': ['{win_1}', '{win_2}'],
             'right_header': 'Challenges', 'right_items': ['{challenge_1}', '{challenge_2}']
            },
            {'type': 'content', 'title': '{quarter} Priorities', 'bullets': [
                '{priority_1}', '{priority_2}', '{priority_3}'
            ]}
        ]
    },
    
    'sales_pitch': {
        'name': 'Sales Pitch Deck',
        'category': 'business',
        'description': 'Product pitch for prospects',
        'theme': 'startup_vibrant',
        'slides': [
            {'type': 'title', 'title': '{product_name}', 'subtitle': '{tagline}'},
            {'type': 'content', 'title': 'The Problem', 'bullets': [
                '{pain_1}', '{pain_2}', '{pain_3}'
            ]},
            {'type': 'content', 'title': 'Our Solution', 'bullets': [
                '{benefit_1}', '{benefit_2}', '{benefit_3}'
            ]},
            {'type': 'content', 'title': 'Pricing', 'bullets': [
                '{price_1}', '{price_2}', '{price_3}'
            ]}
        ]
    },
    
    'investor_pitch': {
        'name': 'Investor Pitch',
        'category': 'business',
        'description': 'Fundraising deck',
        'theme': 'executive_minimal',
        'slides': [
            {'type': 'title', 'title': '{company_name}', 'subtitle': '{tagline}'},
            {'type': 'content', 'title': 'Market Opportunity', 'bullets': [
                'Market size: {market_size}',
                'Growth rate: {growth_rate}',
                'Target segment: {target}'
            ]},
            {'type': 'content', 'title': 'The Ask', 'bullets': [
                'Raising: {raise_amount}',
                'Valuation: {valuation}',
                'Use of funds: {use_1}, {use_2}'
            ]}
        ]
    },
    
    'campaign_review': {
        'name': 'Marketing Campaign Review',
        'category': 'marketing',
        'description': 'Campaign performance analysis',
        'theme': 'marketing_dynamic',
        'slides': [
            {'type': 'title', 'title': 'Campaign Review', 'subtitle': '{campaign_name}'},
            {'type': 'content', 'title': 'Overview', 'bullets': [
                'Objective: {objective}',
                'Budget: {budget}',
                'Duration: {duration}'
            ]},
            {'type': 'content', 'title': 'Results', 'bullets': [
                'Impressions: {impressions}',
                'Engagement: {engagement}',
                'ROI: {roi}'
            ]}
        ]
    },
    
    'product_launch': {
        'name': 'Product Launch',
        'category': 'marketing',
        'description': 'Go-to-market strategy',
        'theme': 'creative_bold',
        'slides': [
            {'type': 'title', 'title': '{product_name}', 'subtitle': 'Product Launch Plan'},
            {'type': 'content', 'title': 'Product Overview', 'bullets': [
                '{description}',
                'Target: {target}',
                'Launch: {date}'
            ]},
            {'type': 'content', 'title': 'GTM Strategy', 'bullets': [
                '{channel_1}', '{channel_2}', '{channel_3}'
            ]}
        ]
    },
    
    'course_overview': {
        'name': 'Course Overview',
        'category': 'education',
        'description': 'Course syllabus introduction',
        'theme': 'education_friendly',
        'slides': [
            {'type': 'title', 'title': '{course_name}', 'subtitle': '{instructor_name}'},
            {'type': 'content', 'title': 'Learning Objectives', 'bullets': [
                '{objective_1}', '{objective_2}', '{objective_3}'
            ]},
            {'type': 'content', 'title': 'Course Topics', 'bullets': [
                '{topic_1}', '{topic_2}', '{topic_3}'
            ]}
        ]
    },
    
    'policy_briefing': {
        'name': 'Policy Briefing',
        'category': 'government',
        'description': 'Policy proposal presentation',
        'theme': 'software_professional',
        'slides': [
            {'type': 'title', 'title': '{policy_name}', 'subtitle': 'Policy Briefing'},
            {'type': 'content', 'title': 'Executive Summary', 'bullets': [
                '{summary_1}', '{summary_2}', '{summary_3}'
            ]},
            {'type': 'content', 'title': 'Key Provisions', 'bullets': [
                '{provision_1}', '{provision_2}', '{provision_3}'
            ]}
        ]
    },
}

def get_template(template_id):
    """Get a template by ID"""
    return TEMPLATE_LIBRARY.get(template_id)

def list_templates():
    """List all templates by category"""
    categories = {}
    for tid, template in TEMPLATE_LIBRARY.items():
        cat = template['category']
        if cat not in categories:
            categories[cat] = []
        categories[cat].append((tid, template['name'], template['description']))
    
    print("\n" + "="*70)
    print("SlideCraft v5 Template Library")
    print("="*70)
    for category, templates in sorted(categories.items()):
        print(f"\n{category.upper()}:")
        for tid, name, desc in templates:
            print(f"  • {name} ({tid})")
            print(f"    {desc}")
    print()

# ============================================================================
# PRESENTATION BUILDER
# ============================================================================

def create_presentation(config):
    """Create v5 presentation"""
    return PresentationBuilderPro(config).build()

# ============================================================================
# TESTS
# ============================================================================

if __name__ == '__main__':
    print("="*70)
    print("🎨 SLIDECRAFT v5.0 - Presentations that inspire, in seconds")
    print("="*70 + "\n")
    
    # Show template library
    list_templates()
    
    # Test: Church sermon
    print("\nTEST: Church Sermon Template")
    print("-"*70)
    r1 = create_presentation({
        'theme': 'church_warmth',
        'slides_content': [
            {'type': 'title', 'title': 'Walking in Faith', 'subtitle': 'Hebrews 11:1-6\nNovember 17, 2025'},
            {'type': 'section', 'title': 'Opening Scripture', 'section_number': '01'},
            {'type': 'content', 'title': 'Context & Background', 'bullets': [
                'Written to Jewish Christians facing persecution',
                'Chapter 11 is the Hall of Faith',
                'Foundation for understanding biblical faith'
            ]},
            {'type': 'content', 'title': 'Practical Application', 'bullets': [
                'This week: Trust God with one specific concern',
                'In relationships: Show faith through actions',
                'Long-term: Build a pattern of faithful living'
            ]}
        ]
    })
    
    print("\nAll tests complete! ✓")
    print("="*70)

